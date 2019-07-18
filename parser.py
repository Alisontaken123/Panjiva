import pandas as pd
import numpy as np
from pptx import *
from pptx.chart.data import ChartData
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
import copy
import geopandas # make world heat map plot

from panjiva import *

class parser:
    def __init__(self, file, china_exports, us_imports, us_imports_12):
        self.file = file
        self.prs = Presentation(self.file)
        self.slide = None 
        self.china_exports = china_exports
        self.us_imports = us_imports
        self.us_imports_12 = us_imports_12

    def parse_table(self, df, table):
        '''parse a pandas dataframe into a powerpoint table
        '''
        # append rows to table
        for i in range(df.shape[0]-1):
            new_row = copy.deepcopy(table._tbl.tr_lst[-1])
            table._tbl.append(new_row)


        for i in range(df.shape[0]):
            for j in range(df.shape[1]):
                table.cell(i+1,j).text = df.iloc[i,j]
                if table.cell(i+1, j).text:
                    font = table.cell(i+1,j).text_frame.paragraphs[0].runs[0].font
                    font.color.rgb = RGBColor(0, 0, 0)
                    font.size = Pt(9)
                    font.name = 'Calibri'
                    font.bold = False
        return table

    def parse_summary_sentence(self, text_frame, input_text):
        '''parse input_text into desired text_frame
        '''
        text_frame.paragraphs[0].text = input_text
        font = text_frame.paragraphs[0].font
        font.color.rgb = RGBColor(0, 0, 0)
        font.size = Pt(9)
        font.name = 'Calibri'
        font.bold = False
        return text_frame

    def parse_slide(self, slide, df, summary_sentence=None):
        '''Given a template slide, data frame and summary sentence, edit the slide in-place. 
        '''
        for shape in slide.shapes:
            if shape.has_table == True:
                table = shape.table 
                table = self.parse_table(df, table)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            if text_frame.paragraphs[0].text == 'Input summary sentence':
                text_frame = self.parse_summary_sentence(text_frame, summary_sentence) 


    def parse_shipment_destinations(self):
        """parse shipment destinations dataframe into slide
        """
        slide = self.prs.slides[2]
        df = shipment_destinations(self.china_exports)
        summary_sentence = exports_summary_sentences(self.china_exports)
        self.parse_slide(slide, df, summary_sentence)

        countries_df = self.china_exports.groupby(['Shipment Destination']).agg(
        {"Value of Goods (USD)": np.sum}).reset_index().sort_values(
        by=['Value of Goods (USD)'], ascending=False)

        # geopandas dataframe
        world = geopandas.read_file(geopandas.datasets.get_path('naturalearth_lowres'))
        world = world.merge(countries_df, how='left', left_on='name', right_on='Shipment Destination')

        # draw geopandas heat world map based on export values
        plot = world.plot(figsize=(5,3), column='Value of Goods (USD)', cmap='Reds', legend=True)
        plot.axis('off')
        fig = plot.get_figure()
        fig.savefig("output.png")

        shapes = slide.shapes
        pictures = shapes.add_picture('output.png', Inches(0.57), Inches(2.06))

        self.prs.save('test.pptx')

    def parse_yearly_exports(self):
        """parse yearly exports dataframe into slide
        """
        slide = self.prs.slides[3]
        df = yearly_exports(self.china_exports)
        summary_sentence = hs_exports_summary_sentence(self.china_exports)
        self.parse_slide(slide, df, summary_sentence)

        for shape in slide.shapes:
            if shape.has_chart == True:
                chart = shape.chart

                chart_data = ChartData()
                years = [str(df['year'][i]) for i in range(5)]
                chart_data.categories = years[0], years[1], years[2], years[3], years[4]
                chart_data.add_series('Total', tuple(df['Total']), number_format='#,##0')
                chart_data.add_series('US', tuple(df['US']), number_format='#,##0')

                chart.replace_data(chart_data)

        self.prs.save('test.pptx')

    def parse_hs_exports(self):
        """parse hs exports dataframe into slide
        """
        slide = self.prs.slides[4]
        df = hs_exports(self.china_exports)
        summary_sentence = hs_exports_summary_sentence(self.china_exports)
        self.parse_slide(slide, df, summary_sentence)

        self.prs.save('test.pptx')

    def parse_yearly_imports(self):
        """parse yearly imports into slide
        """
        slide = self.prs.slides[5]
        df = yearly_imports(self.us_imports)
        summary_sentence = yearly_imports_summary_sentence(self.us_imports)
        self.parse_slide(slide, df, summary_sentence)

        # paste data into graph
        shape = slide.shapes[-2]
        chart = shape.chart
        chart_data = ChartData()
        years = [str(df['year'][i]) for i in range(6)]
        chart_data.categories = years[0],years[1],years[2],years[3],years[4],years[5]
        chart_data.add_series('Number of Shipments', tuple(df['Number of Shipments']), number_format='#,##0')
        chart_data.add_series('Number of Containers', tuple(df['Number of Containers']), number_format='#,##0')
        chart.replace_data(chart_data)

        # paste data into graph
        df = monthly_imports(self.us_imports_12)
        shape = slide.shapes[-1]
        chart = shape.chart
        chart_data = ChartData()
        months = [str(df['month'][i]) for i in range(13)]
        chart_data.categories = months[0],months[1],months[2], months[3], months[4],months[5],months[6],months[7],months[8],months[9],months[10],months[11],months[12]
        chart_data.add_series('Number of Shipments', tuple(df['Number of Shipments']), number_format='#,##0')
        chart_data.add_series('Number of Containers', tuple(df['Number of Containers']), number_format='#,##0')

        chart.replace_data(chart_data)

        self.prs.save('test.pptx')

    def parse_hs_imports(self):
        """parse hs imports dataframe into slide
        """
        slide = self.prs.slides[6]
        df = hs_imports_merge_12(self.us_imports,self.us_imports_12)
        summary_sentence = hs_imports_summary_sentence(self.us_imports)
        self.parse_slide(slide, df, summary_sentence)
        self.prs.save('test.pptx')

    def parse_consignees_imports(self):
        """parse consignees imports dataframe into slide
        """
        slide = self.prs.slides[7]
        df = consignees_imports(self.us_imports)
        summary_sentence = consignees_imports_summary_sentence(self.us_imports)
        self.parse_slide(slide, df, summary_sentence)
        self.prs.save('test.pptx')

    def parse_consignees_imports_12(self):
        """parse consignees imports dataframe into slide
        """
        slide = self.prs.slides[8]
        df = consignees_imports_12(self.us_imports_12)
        summary_sentence = consignees_imports_12_summary_sentence(self.us_imports)
        self.parse_slide(slide, df, summary_sentence)
        self.prs.save('test.pptx') 

    def parse_recent_shipments(self):
        """parse recent shipments dataframe into slide
        """
        slide = self.prs.slides[9]
        df = recent_shipments(self.us_imports)
        self.parse_slide(slide, df)
        self.prs.save('test.pptx')



